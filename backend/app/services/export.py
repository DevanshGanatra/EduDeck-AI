import os
import uuid
from typing import List
from pptx import Presentation
from pptx.util import Inches, Pt
from app.schemas.generation import SlideContent
from app.services.storage import StorageService

class ExportService:
    def __init__(self, storage_service: StorageService):
        self.storage_service = storage_service

    async def create_pptx(self, slides: List[SlideContent]) -> str:
        prs = Presentation()
        
        for slide_data in slides:
            # 1 (Title and Content layout)
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = slide_data.title
                
            # Content
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = "\n".join(f"• {bp}" for bp in slide_data.bullet_points)
            
            # Add References if any
            if slide_data.references:
                ref_text = "References: " + ", ".join(slide_data.references)
                # Add a small textbox at the bottom
                left = Inches(0.5)
                top = Inches(6.8)
                width = Inches(9)
                height = Inches(0.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                p = txBox.text_frame.add_paragraph()
                p.text = ref_text
                p.font.size = Pt(10)
                
            # Speaker Notes
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data.speaker_notes

        filename = f"presentation_{uuid.uuid4()}.pptx"
        # We assume local storage for MVP, save directly
        # In a real impl, we'd save to a BytesIO buffer and use self.storage_service.save_upload()
        # For simplicity of MVP local saving:
        os.makedirs("data/generated", exist_ok=True)
        file_path = os.path.join("data/generated", filename)
        prs.save(file_path)
        
        return f"/download/{filename}" # Mock URL
